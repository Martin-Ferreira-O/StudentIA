# app.py
import os
from flask import Flask, request, render_template, jsonify, session
from pptx import Presentation
import whisper
from datetime import datetime
import db
from models import Usuario, Resumen
from database_manager import UsuarioManager, ResumenManager
from flask.cli import with_appcontext
import click

from ia import Resumen as ResumenIA

# Función para crear tablas de BD
def crear_tablas_bd():
    """Crear tablas de la base de datos"""
    try:
        # Crear la base de datos si no existe
        from sqlalchemy import create_engine
        from sqlalchemy_utils import database_exists, create_database
        
        # Obtener la URI de la base de datos del objeto app
        engine = create_engine(app.config['SQLALCHEMY_DATABASE_URI'])
        
        if not database_exists(engine.url):
            print(f"Creando base de datos: {engine.url.database}")
            create_database(engine.url)
            print("Base de datos creada correctamente")
        
        # Crear las tablas
        print("Creando tablas...")
        db.db.create_all()
        print("Tablas creadas correctamente")
        
    except Exception as e:
        print(f"Error al crear tablas: {e}")
        print("Asegúrate de que PostgreSQL esté instalado y ejecutándose.")
        print("Verifica las credenciales en el archivo .env")

app = Flask(__name__)
app.secret_key = os.getenv('SECRET_KEY', 'dev_key_very_secret')

# Inicializar la base de datos
db.init_app(app)

# Registrar comando para crear las tablas
@click.command('crear-tablas')
@with_appcontext
def crear_tablas_comando():
    """Comando para crear tablas de la base de datos desde la CLI."""
    crear_tablas_bd()
    click.echo('Tablas de base de datos creadas.')

app.cli.add_command(crear_tablas_comando)

# Inicializar la IA para resúmenes
resumen_ia = ResumenIA()

# Función para guardar texto en archivo TXT
def guardar_texto_txt(texto, prefijo="texto"):
    # Crear directorio de salida si no existe
    if not os.path.exists("output"):
        os.makedirs("output")
    
    # Generar nombre de archivo con fecha unica
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    nombre_archivo = f"output/{prefijo}_{timestamp}.txt"
    
    # Guardar el texto en el archivo
    with open(nombre_archivo, "w", encoding="utf-8") as f:
        f.write(texto)
    
    resumen_ia.read_files(nombre_archivo)

    return nombre_archivo

# Función para extraer texto de un archivo .pptx
def extraer_texto_pptx(archivo_pptx):
    presentacion = Presentation(archivo_pptx)
    texto_completo = ""
    for diapositiva in presentacion.slides:
        for forma in diapositiva.shapes:
            if hasattr(forma, "text"):
                texto_completo += forma.text + "\n"
    
    # Guardar el texto extraído en un archivo TXT
    archivo_txt = guardar_texto_txt(texto_completo, "pptx")
    print(f"Texto extraído guardado en: {archivo_txt}")
    
    return texto_completo

# Función para transcribir audio usando Whisper
def transcribir_audio(audio_path):
    model = whisper.load_model("small") # large, medium, small, base
    result = model.transcribe(audio_path, language="es")
    texto_transcrito = result["text"]
    
    # Guardar la transcripción en un archivo TXT
    archivo_txt = guardar_texto_txt(texto_transcrito, "audio")
    print(f"Transcripción guardada en: {archivo_txt}")
    
    return texto_transcrito


# Rutas de la aplicación
@app.route("/", methods=["GET", "POST"])
def index():
    if request.method == "POST":
        archivo = request.files["archivo"]
        tipo_archivo = request.form.get("tipo_archivo")
        salidas = request.form.getlist("salida")
                
        if tipo_archivo == "pptx":
            # Guardar archivo temporalmente
            ruta_temporal = os.path.join("temp", archivo.filename)
            archivo.save(ruta_temporal)

            # Extraer texto del archivo .pptx
            texto = extraer_texto_pptx(ruta_temporal)
            os.remove(ruta_temporal)  # Eliminar archivo temporal

        elif tipo_archivo == "audio":
            # Guardar archivo temporalmente
            ruta_temporal = os.path.join("temp", archivo.filename)
            archivo.save(ruta_temporal)

            # Transcribir audio
            texto = transcribir_audio(ruta_temporal)
            os.remove(ruta_temporal)  # Eliminar archivo temporal

        else:
            return jsonify({"error": "Tipo de archivo no soportado"}), 400

        resultado = {"texto": texto}
        
        if "resumen" in salidas:
            contenido_resumen = resumen_ia.send_resume()
            resultado["resumen"] = contenido_resumen
        if "mapa" in salidas:
            mapa = resumen_ia.generate_concept_map()
            print("\n--- CÓDIGO MERMAID GENERADO ---\n", mapa, "\n------------------------------\n")
            resultado["mapa_conceptual"] = mapa
        if "quiz" in salidas:
            quiz = resumen_ia.generate_quiz()
            resultado["quiz"] = quiz
            
        # Si el usuario está logueado, guardar el resumen en la base de datos
        usuario_id = session.get('usuario_id')
        if usuario_id:
            try:
                titulo = archivo.filename
                ResumenManager.crear_resumen(
                    titulo=titulo,
                    texto_original=texto,
                    tipo_archivo=tipo_archivo,
                    usuario_id=usuario_id,
                    contenido_resumen=resultado.get('resumen'),
                    contenido_mapa=resultado.get('mapa_conceptual'),
                    contenido_quiz=resultado.get('quiz')
                )
            except Exception as e:
                print(f"Error al guardar resumen: {e}")

        return jsonify(resultado)

    return render_template("index.html")

@app.route("/login", methods=["GET", "POST"])
def login():
    error_message = None
    
    if request.method == "POST":
        email = request.form.get("email")
        password = request.form.get("password")
        
        # Buscar usuario por email
        usuario = UsuarioManager.obtener_usuario_por_email(email)
        
        if usuario and usuario.check_password(password):
            # Iniciar sesión
            session['usuario_id'] = usuario.id
            session['usuario_nombre'] = usuario.nombre
            return jsonify({"success": True})
        else:
            error_message = "Correo o contraseña incorrectos"
            return jsonify({"success": False, "error": error_message})
    
    return render_template("login.html", error=error_message)

@app.route("/registro", methods=["GET", "POST"])
def registro():
    error_message = None
    
    if request.method == "POST":
        nombre = request.form.get("nombre")
        email = request.form.get("email")
        password = request.form.get("password")
        fecha_nacimiento = request.form.get("fecha_nacimiento")
        
        if not nombre or not email or not password:
            error_message = "Todos los campos son obligatorios"
            return jsonify({"success": False, "error": error_message})
        
        # Convertir fecha_nacimiento a objeto datetime si existe
        if fecha_nacimiento:
            try:
                fecha_nacimiento = datetime.strptime(fecha_nacimiento, "%Y-%m-%d").date()
            except ValueError:
                error_message = "Formato de fecha inválido"
                return jsonify({"success": False, "error": error_message})
        
        # Crear usuario
        usuario, error = UsuarioManager.crear_usuario(nombre, email, password, fecha_nacimiento)
        
        if error:
            return jsonify({"success": False, "error": error})
        
        # Iniciar sesión automáticamente
        session['usuario_id'] = usuario.id
        session['usuario_nombre'] = usuario.nombre
        
        return jsonify({"success": True})
    
    return render_template("registro.html", error=error_message)

@app.route("/logout")
def logout():
    # Cerrar sesión
    session.pop('usuario_id', None)
    session.pop('usuario_nombre', None)
    return jsonify({"success": True})

@app.route("/mis-resumenes")
def mis_resumenes():
    # Verificar si el usuario está logueado
    usuario_id = session.get('usuario_id')
    if not usuario_id:
        return render_template("login.html", error="Debes iniciar sesión para ver tus resúmenes")
    
    # Obtener todos los resúmenes del usuario
    resumenes = ResumenManager.obtener_resumenes_por_usuario(usuario_id)
    
    return render_template("mis_resumenes.html", resumenes=resumenes)

@app.route("/resumen/<int:resumen_id>")
def ver_resumen(resumen_id):
    # Verificar si el usuario está logueado
    usuario_id = session.get('usuario_id')
    if not usuario_id:
        return render_template("login.html", error="Debes iniciar sesión para ver este resumen")
    
    # Obtener el resumen
    resumen = ResumenManager.obtener_resumen_por_id(resumen_id)
    
    if not resumen or resumen.usuario_id != usuario_id:
        return render_template("error.html", error="Resumen no encontrado o no tienes permisos para verlo")
    
    return render_template("ver_resumen.html", resumen=resumen)

@app.route("/resumen/<int:resumen_id>/eliminar", methods=["POST"])
def eliminar_resumen(resumen_id):
    # Verificar si el usuario está logueado
    usuario_id = session.get('usuario_id')
    if not usuario_id:
        return jsonify({"success": False, "error": "Debes iniciar sesión para eliminar este resumen"})
    
    # Obtener el resumen
    resumen = ResumenManager.obtener_resumen_por_id(resumen_id)
    
    if not resumen:
        return jsonify({"success": False, "error": "Resumen no encontrado"})
    
    if resumen.usuario_id != usuario_id:
        return jsonify({"success": False, "error": "No tienes permisos para eliminar este resumen"})
    
    # Eliminar el resumen
    exito, error = ResumenManager.eliminar_resumen(resumen_id)
    
    if exito:
        return jsonify({"success": True})
    else:
        return jsonify({"success": False, "error": error})

# Manejador de error 404 - Página no encontrada
@app.errorhandler(404)
def page_not_found(e):
    return render_template('error404.html'), 404

if __name__ == "__main__":
    if not os.path.exists("temp"):
        os.makedirs("temp")
    if not os.path.exists("output"):
        os.makedirs("output")
    
    # Crear tablas de BD al iniciar la aplicación
    with app.app_context():
        crear_tablas_bd()
    
    print("Iniciando servidor... en http://localhost:3000")
    app.run(debug=True, port=3000)
